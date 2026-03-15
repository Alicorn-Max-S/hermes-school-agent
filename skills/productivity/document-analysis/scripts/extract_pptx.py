#!/usr/bin/env python3
"""Extract text from PowerPoint files using python-pptx (FREE, local).

Usage:
    python3 extract_pptx.py presentation.pptx
"""

import json
import sys
from pathlib import Path


def extract(input_path):
    """Extract text from a PPTX file."""
    path = Path(input_path)
    if not path.exists():
        return {"success": False, "error": f"File not found: {input_path}"}

    try:
        from pptx import Presentation
    except ImportError:
        return {
            "success": False,
            "error": "python-pptx not installed. Run: pip install python-pptx",
        }

    try:
        prs = Presentation(str(path))
    except Exception as e:
        return {"success": False, "error": f"Failed to open PPTX: {e}"}

    slides_content = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_lines = [f"--- Slide {slide_num} ---"]

        # Extract title if present
        if slide.shapes.title and slide.shapes.title.text.strip():
            slide_lines.append(f"Title: {slide.shapes.title.text.strip()}")

        # Extract text from all shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_lines.append(text)

            # Extract table content
            if shape.has_table:
                table = shape.table
                slide_lines.append("")
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_lines.append("| " + " | ".join(cells) + " |")

        # Extract notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_lines.append(f"\nSpeaker Notes: {notes}")

        slides_content.append("\n".join(slide_lines))

    content = "\n\n".join(slides_content)

    return {
        "success": True,
        "content": content,
        "slides": len(prs.slides),
    }


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(json.dumps({"error": "Usage: extract_pptx.py <file>"}))
        sys.exit(1)

    result = extract(sys.argv[1])
    print(json.dumps(result, indent=2, ensure_ascii=False))
